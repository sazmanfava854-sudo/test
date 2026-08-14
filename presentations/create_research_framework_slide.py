#!/usr/bin/env python3
"""Generate 'چارچوب کلی پژوهش' horizontal timeline slide (reference style)."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

# Colors from reference
BLUE = RGBColor(0x25, 0x63, 0xEB)
BLUE_DARK = RGBColor(0x1D, 0x4E, 0xD8)
BLUE_LIGHT = RGBColor(0x3B, 0x82, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK = RGBColor(0x1E, 0x29, 0x3B)
TEXT_MUTED = RGBColor(0x64, 0x74, 0x8B)
ARROW_GRAY = RGBColor(0x47, 0x55, 0x69)
SIDEBAR_BLUE = RGBColor(0x25, 0x63, 0xEB)

STEPS = [
    {
        "title": "خوشه‌بندی",
        "subtitle": "گروه‌بندی فناوری‌های مشابه",
        "large": True,
    },
    {
        "title": "AHP",
        "subtitle": "تعیین اهمیت نسبی معیارها",
        "large": True,
    },
    {
        "title": "تعدیل پویا",
        "subtitle": "تطبیق وزن‌ها با سناریوی پروژه",
        "large": True,
    },
    {
        "title": "TOPSIS",
        "subtitle": "رتبه‌بندی نهایی در خوشه منتخب",
        "large": True,
    },
    {
        "title": "فناوری منتخب",
        "subtitle": "",
        "large": False,
    },
]

# Sidebar icon labels (decorative simple shapes)
SIDEBAR_ICONS = ["book", "map", "lightbulb", "chart", "people"]
ACTIVE_ICON = 2  # lightbulb = روش پیشنهادی


def set_rtl(paragraph):
    p_pr = paragraph._p.get_or_add_pPr()
    p_pr.set(qn("a:rtl"), "1")


def fill_shape(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_textbox(slide, left, top, width, height, text, *, font_size=14, bold=False,
                color=TEXT_DARK, align=PP_ALIGN.CENTER, rtl=True, font_name="Tahoma"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    if rtl:
        set_rtl(p)
    if p.runs:
        run = p.runs[0]
    else:
        run = p.add_run()
        run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return box


def add_rounded_rect(slide, left, top, width, height, fill, border=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
  # round corners more
    if shape.adjustments:
        shape.adjustments[0] = 0.12
    return shape


def add_arrow_shape(slide, left, top, width, height):
    """Thick grey arrow pointing left (RTL flow)."""
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = ARROW_GRAY
    arrow.line.fill.background()
    arrow.rotation = 180.0  # point left
    return arrow


def draw_sidebar_icon(slide, cx, cy, icon_type, active=False):
    size = Inches(0.28)
    if active:
        ring = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, cx - Inches(0.22), cy - Inches(0.22),
            Inches(0.44), Inches(0.44),
        )
        ring.fill.solid()
        ring.fill.fore_color.rgb = WHITE
        ring.line.fill.background()

    s = size
    if icon_type == "book":
        r = add_rounded_rect(slide, cx - s / 2, cy - s / 2, s, s * 0.8, WHITE if active else RGBColor(0xDB, 0xEA, 0xFE))
        add_textbox(slide, cx - s / 2, cy - s / 2, s, s * 0.8, "📖", font_size=10, rtl=False)
    elif icon_type == "map":
        shape = slide.shapes.add_shape(MSO_SHAPE.PENTAGON, cx - s / 2, cy - s / 2, s, s)
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE if active else RGBColor(0xDB, 0xEA, 0xFE)
        shape.line.fill.background()
    elif icon_type == "lightbulb":
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - s / 2, cy - s / 2, s, s)
        shape.fill.solid()
        shape.fill.fore_color.rgb = BLUE if active else RGBColor(0xDB, 0xEA, 0xFE)
        shape.line.fill.background()
        if active:
            add_textbox(slide, cx - s / 2, cy - s / 2, s, s, "💡", font_size=11, rtl=False)
    elif icon_type == "chart":
        shape = slide.shapes.add_shape(MSO_SHAPE.ACTION_BUTTON_DOCUMENT, cx - s / 2, cy - s / 2, s, s)
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE if active else RGBColor(0xDB, 0xEA, 0xFE)
        shape.line.fill.background()
    else:  # people
        shape = slide.shapes.add_shape(MSO_SHAPE.ACTION_BUTTON_HOME, cx - s / 2, cy - s / 2, s, s)
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE if active else RGBColor(0xDB, 0xEA, 0xFE)
        shape.line.fill.background()


def add_step_box(slide, left, top, width, height, step):
    box = add_rounded_rect(slide, left, top, width, height, BLUE)

    if step["subtitle"]:
        add_textbox(
            slide, left + Inches(0.08), top + Inches(0.18), width - Inches(0.16), Inches(0.38),
            step["title"], font_size=15 if step["large"] else 12, bold=True, color=WHITE,
        )
        add_textbox(
            slide, left + Inches(0.08), top + Inches(0.55), width - Inches(0.16), Inches(0.45),
            step["subtitle"], font_size=9, color=RGBColor(0xDB, 0xEA, 0xFE),
        )
    else:
        add_textbox(
            slide, left + Inches(0.06), top + Inches(0.2), width - Inches(0.12), Inches(0.5),
            step["title"], font_size=11, bold=True, color=WHITE,
        )


def build_slide(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # White background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    fill_shape(bg, BG)
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

    # Right sidebar
    sidebar_w = Inches(0.55)
    sidebar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        prs.slide_width - sidebar_w, 0, sidebar_w, prs.slide_height,
    )
    fill_shape(sidebar, SIDEBAR_BLUE)

    # Sidebar icons
    icon_x = prs.slide_width - sidebar_w / 2
    icon_positions = [Inches(0.9), Inches(1.7), Inches(2.55), Inches(3.4), Inches(4.25)]
    for i, (icon, y) in enumerate(zip(SIDEBAR_ICONS, icon_positions)):
        draw_sidebar_icon(slide, icon_x, y, icon, active=(i == ACTIVE_ICON))

    # Title
    add_textbox(
        slide, Inches(0.5), Inches(0.35), Inches(8.5), Inches(0.55),
        "چارچوب کلی پژوهش",
        font_size=26, bold=True, color=TEXT_DARK, align=PP_ALIGN.RIGHT,
    )

    # Timeline area
    timeline_y = Inches(2.35)
    large_w = Inches(1.55)
    large_h = Inches(1.05)
    small_w = Inches(1.15)
    small_h = Inches(0.75)
    arrow_w = Inches(0.38)
    arrow_h = Inches(0.22)
    gap = Inches(0.08)

    # Calculate total width and start position (RTL: first step on right)
    total_w = 4 * (large_w + arrow_w + gap) + small_w
    start_x = prs.slide_width - sidebar_w - Inches(0.45) - total_w

    x = start_x + total_w - small_w  # place from right

    # Draw steps RTL: خوشه‌بندی (rightmost) → فناوری منتخب (leftmost)
    positions = []
    for i, step in enumerate(STEPS):
        w = large_w if step["large"] else small_w
        h = large_h if step["large"] else small_h
        y = timeline_y if step["large"] else timeline_y + (large_h - small_h) / 2
        positions.append((x, y, w, h, step))
        x -= w + arrow_w + gap

    # Reverse to draw left-to-right in z-order but positions are RTL
    positions.reverse()
    # Recalculate positions properly RTL
    x_right = prs.slide_width - sidebar_w - Inches(0.5)
    ordered = []
    cx = x_right
    for step in STEPS:
        w = large_w if step["large"] else small_w
        h = large_h if step["large"] else small_h
        y = timeline_y if step["large"] else timeline_y + (large_h - small_h) / 2
        left = cx - w
        ordered.append((left, y, w, h, step))
        cx = left - arrow_w - gap

    # Draw arrows between consecutive steps (pointing left / RTL)
    for i in range(len(ordered) - 1):
        left1, y1, w1, h1, _ = ordered[i]       # right-side box
        left2, y2, w2, h2, _ = ordered[i + 1]   # left-side box
        arrow_left = left2 + w2
        arrow_width = left1 - arrow_left
        if arrow_width > Inches(0.1):
            arrow_top = timeline_y + large_h / 2 - arrow_h / 2
            add_arrow_shape(slide, arrow_left, arrow_top, arrow_width, arrow_h)

    # Draw boxes
    for left, y, w, h, step in ordered:
        add_step_box(slide, left, y, w, h, step)

    # Footer
    add_textbox(
        slide, Inches(0.5), Inches(4.85), Inches(2.5), Inches(0.35),
        "روش پیشنهادی", font_size=12, bold=True, color=BLUE, align=PP_ALIGN.RIGHT,
    )

    # Page nav mockup
    nav_y = Inches(4.88)
    prev_c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.0), nav_y, Inches(0.28), Inches(0.28))
    fill_shape(prev_c, BLUE)
    add_textbox(slide, Inches(3.35), nav_y - Inches(0.02), Inches(0.35), Inches(0.3), "15", font_size=11, bold=True, color=TEXT_DARK, rtl=False)
    next_c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.75), nav_y, Inches(0.28), Inches(0.28))
    fill_shape(next_c, BLUE)


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    build_slide(prs)
    output = "/workspace/presentations/چارچوب_کلی_پژوهش_تایم‌لاین.pptx"
    prs.save(output)
    print(f"Saved: {output}")


if __name__ == "__main__":
    main()
