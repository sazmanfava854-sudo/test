#!/usr/bin/env python3
"""Generate 4-phase system timeline slide (central-axis graphic style)."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# Palette
NAVY = RGBColor(0x1E, 0x3A, 0x5F)
NAVY_DARK = RGBColor(0x15, 0x2A, 0x45)
NAVY_LIGHT = RGBColor(0x2B, 0x5C, 0x8A)
TEAL = RGBColor(0x00, 0x96, 0x88)
TEAL_DARK = RGBColor(0x00, 0x7A, 0x6E)
TEAL_BG = RGBColor(0xE0, 0xF2, 0xF1)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
GREEN_BG = RGBColor(0xE8, 0xF5, 0xE9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG = RGBColor(0xF7, 0xF9, 0xFC)
TEXT = RGBColor(0x2D, 0x34, 0x3E)
MUTED = RGBColor(0x6B, 0x72, 0x80)
BORDER = RGBColor(0xD1, 0xD9, 0xE6)

PHASES = [
    {
        "num": "۱",
        "title": "خوشه‌بندی فناوری‌ها",
        "desc": "گروه‌بندی بر اساس ویژگی‌های فنی و اقتصادی",
        "pos": "above",
        "node_color": NAVY_LIGHT,
        "label": "شروع",
        "label_pos": "below_node",
        "final": False,
        "highlight": False,
    },
    {
        "num": "۲",
        "title": "محاسبه وزن پایه (AHP)",
        "desc": "بردار وزن معیارها از قضاوت خبرگان",
        "pos": "below",
        "node_color": NAVY_LIGHT,
        "label": "",
        "final": False,
        "highlight": False,
    },
    {
        "num": "۳",
        "title": "تعدیل پویای وزن",
        "desc": "تبدیل وزن پایه به وزن تطبیقی — نوآوری",
        "pos": "above",
        "node_color": TEAL,
        "label": "نوآوری",
        "label_pos": "above_node",
        "final": False,
        "highlight": True,
    },
    {
        "num": "۴",
        "title": "رتبه‌بندی (TOPSIS)",
        "desc": "انتخاب فناوری برتر در خوشه منتخب",
        "pos": "below",
        "node_color": NAVY,
        "label": "پایان",
        "label_pos": "below_node",
        "final": True,
        "highlight": False,
    },
]


def set_rtl(paragraph):
    p_pr = paragraph._p.get_or_add_pPr()
    p_pr.set(qn("a:rtl"), "1")


def fill_shape(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_textbox(slide, left, top, width, height, text, *, font_size=12, bold=False,
                color=TEXT, align=PP_ALIGN.CENTER, rtl=True, font_name="Tahoma"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    if rtl:
        set_rtl(p)
    run = p.runs[0] if p.runs else p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return box


def add_rounded_rect(slide, left, top, width, height, fill, border=None, border_w=Pt(1.5)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border:
        shape.line.color.rgb = border
        shape.line.width = border_w
    else:
        shape.line.fill.background()
    if shape.adjustments:
        shape.adjustments[0] = 0.1
    return shape


def add_node(slide, cx, cy, radius, phase):
    outer = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, cx - radius - Inches(0.05), cy - radius - Inches(0.05),
        (radius + Inches(0.05)) * 2, (radius + Inches(0.05)) * 2,
    )
    outer.fill.solid()
    outer.fill.fore_color.rgb = WHITE
    outer.line.color.rgb = phase["node_color"]
    outer.line.width = Pt(3)

    inner = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - radius, cy - radius, radius * 2, radius * 2)
    inner.fill.solid()
    inner.fill.fore_color.rgb = phase["node_color"]
    inner.line.fill.background()

    inner.text_frame.clear()
    p = inner.text_frame.paragraphs[0]
    p.text = phase["num"]
    p.alignment = PP_ALIGN.CENTER
    inner.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    run = p.runs[0]
    run.font.size = Pt(15)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = "Tahoma"


def add_phase_card(slide, cx, line_y, radius, phase, card_w, card_h):
    is_above = phase["pos"] == "above"
    is_final = phase["final"]
    is_highlight = phase["highlight"]

    left = cx - card_w / 2
    if is_above:
        top = line_y - radius - Inches(0.2) - card_h
        stem_y1 = top + card_h
        stem_y2 = line_y - radius - Inches(0.04)
    else:
        top = line_y + radius + Inches(0.2)
        stem_y1 = line_y + radius + Inches(0.04)
        stem_y2 = top

    card_fill = NAVY if is_final else WHITE
    border = TEAL if is_highlight else (NAVY if is_final else BORDER)
    title_color = WHITE if is_final else TEXT
    desc_color = RGBColor(0xBB, 0xDE, 0xFB) if is_final else MUTED

    add_rounded_rect(slide, left, top, card_w, card_h, card_fill, border, Pt(2 if is_highlight or is_final else 1.2))

    add_textbox(slide, left + Inches(0.1), top + Inches(0.12), card_w - Inches(0.2), Inches(0.32),
                phase["title"], font_size=11, bold=True, color=title_color)
    add_textbox(slide, left + Inches(0.1), top + Inches(0.44), card_w - Inches(0.2), Inches(0.5),
                phase["desc"], font_size=9, color=desc_color)

    if is_highlight:
        badge = add_rounded_rect(slide, left + Inches(0.15), top + card_h - Inches(0.3),
                                 card_w - Inches(0.3), Inches(0.22), GREEN_BG, GREEN, Pt(1))
        badge.text_frame.clear()
        p = badge.text_frame.paragraphs[0]
        p.text = "تعدیل نمایی"
        p.alignment = PP_ALIGN.CENTER
        set_rtl(p)
        run = p.runs[0]
        run.font.size = Pt(8)
        run.font.bold = True
        run.font.color.rgb = GREEN
        run.font.name = "Tahoma"

    # Stem
    stem = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx - Pt(1.2), min(stem_y1, stem_y2),
                                  Pt(2.4), abs(stem_y2 - stem_y1))
    stem.fill.solid()
    stem.fill.fore_color.rgb = phase["node_color"]
    stem.line.fill.background()

    # Phase label near node
    if phase.get("label"):
        label = phase["label"]
        if phase.get("label_pos") == "above_node":
            ly = line_y - radius - Inches(0.38)
        else:
            ly = line_y + radius + Inches(0.12)
        add_textbox(slide, cx - Inches(0.45), ly, Inches(0.9), Inches(0.22),
                    label, font_size=8, bold=True, color=phase["node_color"])


def build_slide(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    fill_shape(bg, BG)
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.1))
    fill_shape(top_bar, NAVY)

    add_textbox(slide, Inches(0.45), Inches(0.28), Inches(9.1), Inches(0.5),
                "چهار فاز اصلی سامانه", font_size=26, bold=True, color=NAVY, align=PP_ALIGN.RIGHT)
    add_textbox(slide, Inches(0.45), Inches(0.78), Inches(9.1), Inches(0.32),
                "از خوشه‌بندی تا رتبه‌بندی نهایی فناوری‌ها", font_size=12, color=MUTED, align=PP_ALIGN.RIGHT)

    line_y = Inches(3.05)
    line_left = Inches(0.9)
    line_right = Inches(9.1)
    line_w = line_right - line_left

    # Base axis (dark blue)
    base = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, line_left, line_y - Inches(0.05), line_w, Inches(0.1))
    fill_shape(base, NAVY)

    # White accent line on top
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, line_left, line_y - Inches(0.07), line_w, Inches(0.025))
    fill_shape(accent, WHITE)

    # Teal section on axis (phase 3 area — middle-right segment)
    n = len(PHASES)
    spacing = line_w / (n - 1)
    positions = [line_right - i * spacing for i in range(n)]

    # Teal segment between node 2 and node 3 (indices 1 and 2)
    teal_left = positions[2] - spacing * 0.45
    teal_right = positions[1] + spacing * 0.45
    teal_seg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, teal_left, line_y - Inches(0.05), teal_right - teal_left, Inches(0.1)
    )
    fill_shape(teal_seg, TEAL)

    radius = Inches(0.3)
    card_w = Inches(2.05)
    card_h = Inches(0.95)

    for phase, cx in zip(PHASES, positions):
        add_node(slide, cx, line_y, radius, phase)
        add_phase_card(slide, cx, line_y, radius, phase, card_w, card_h)

    # Bottom legend strip
    legend = add_rounded_rect(slide, Inches(0.45), Inches(4.55), Inches(9.1), Inches(0.72), WHITE, BORDER)
    legends = [
        (NAVY_LIGHT, "فاز پردازش"),
        (TEAL, "فاز نوآورانه (تعدیل پویا)"),
        (NAVY, "خروجی نهایی"),
    ]
    lx = Inches(1.2)
    for color, label in legends:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, lx, Inches(4.78), Inches(0.14), Inches(0.14))
        fill_shape(dot, color)
        add_textbox(slide, lx + Inches(0.2), Inches(4.72), Inches(2.2), Inches(0.28), label, font_size=9, color=TEXT)
        lx += Inches(2.8)

    add_textbox(slide, Inches(0.45), Inches(5.12), Inches(9.1), Inches(0.22),
                "روش پیشنهادی  |  AHP + تعدیل پویا + TOPSIS",
                font_size=9, color=MUTED, align=PP_ALIGN.LEFT)


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    build_slide(prs)
    output = "/workspace/presentations/چهار_فاز_اصلی_سامانه_تایم‌لاین.pptx"
    prs.save(output)
    print(f"Saved: {output}")


if __name__ == "__main__":
    main()
