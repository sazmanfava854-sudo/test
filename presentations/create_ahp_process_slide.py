#!/usr/bin/env python3
"""Generate a Persian AHP expert-opinion aggregation process flowchart slide."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

# Palette (same as consistency slide)
NAVY = RGBColor(0x1E, 0x3A, 0x5F)
NAVY_LIGHT = RGBColor(0x2B, 0x5C, 0x8A)
ACCENT = RGBColor(0x00, 0x96, 0x88)
TEAL_BG = RGBColor(0xE0, 0xF2, 0xF1)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
GREEN_BG = RGBColor(0xE8, 0xF5, 0xE9)
GREEN_DARK = RGBColor(0x1B, 0x5E, 0x20)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG = RGBColor(0xF7, 0xF9, 0xFC)
TEXT = RGBColor(0x2D, 0x34, 0x3E)
MUTED = RGBColor(0x6B, 0x72, 0x80)
BORDER = RGBColor(0xD1, 0xD9, 0xE6)
RED_SOFT = RGBColor(0xC6, 0x28, 0x28)
RED_BG = RGBColor(0xFF, 0xEB, 0xEE)
AMBER = RGBColor(0xF5, 0x7C, 0x00)
AMBER_BG = RGBColor(0xFF, 0xF3, 0xE0)

STEPS = [
    {
        "num": "۱",
        "title": "گردآوری نظرات خبرگان",
        "desc": "جمع‌آوری ماتریس‌های مقایسه زوجی از خبرگان",
        "fill": WHITE,
        "border": BORDER,
        "accent": NAVY_LIGHT,
    },
    {
        "num": "۲",
        "title": "مقایسه‌های زوجی معیارها",
        "desc": "انجام با مقیاس ساعتی AHP (۱ تا ۹)",
        "fill": WHITE,
        "border": BORDER,
        "accent": NAVY_LIGHT,
    },
    {
        "num": "۳",
        "title": "بررسی سازگاری هر خبره",
        "desc": "محاسبه CR برای ماتریس هر خبره به‌صورت جداگانه",
        "fill": TEAL_BG,
        "border": ACCENT,
        "accent": ACCENT,
    },
    {
        "num": "۴",
        "title": "پذیرش ماتریس‌های سازگار",
        "desc": "فقط ماتریس‌هایی با CR < ۰٫۱ پذیرفته می‌شوند",
        "fill": GREEN_BG,
        "border": GREEN,
        "accent": GREEN,
        "is_filter": True,
    },
    {
        "num": "۵",
        "title": "تجمیع با میانگین هندسی",
        "desc": "ادغام نظرات پذیرفته‌شده خبرگان",
        "fill": WHITE,
        "border": BORDER,
        "accent": NAVY_LIGHT,
    },
    {
        "num": "۶",
        "title": "محاسبه وزن نهایی معیارها",
        "desc": "ماتریس تجمیع‌شده مبنای بردار وزن نهایی است",
        "fill": NAVY,
        "border": NAVY,
        "accent": ACCENT,
        "final": True,
    },
]


def set_rtl(paragraph):
    p_pr = paragraph._p.get_or_add_pPr()
    p_pr.set(qn("a:rtl"), "1")


def fill_shape(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_textbox(slide, left, top, width, height, text, *, font_size=12, bold=False,
                color=TEXT, align=PP_ALIGN.RIGHT, rtl=True, font_name="Tahoma"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    if rtl:
        set_rtl(p)
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return box


def add_rounded_rect(slide, left, top, width, height, fill, border=None, border_width=Pt(1.5)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border:
        shape.line.color.rgb = border
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape


def add_arrow_down(slide, x_center, y_top, length):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, x_center, y_top, x_center, y_top + length
    )
    conn.line.color.rgb = ACCENT
    conn.line.width = Pt(2.5)
    # arrowhead via line end
    ln = conn.line._ln
    tail = parse_xml(
        '<a:tailEnd xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" type="triangle"/>'
    )
    ln.append(tail)
    return conn


def add_step_box(slide, left, top, width, height, step):
    is_final = step.get("final", False)
    is_filter = step.get("is_filter", False)

    box = add_rounded_rect(slide, left, top, width, height, step["fill"], step["border"], Pt(2 if is_filter else 1.5))

    # Number badge
    badge_size = Inches(0.38)
    badge_left = left + width - badge_size - Inches(0.12)
    badge_top = top + Inches(0.1)
    badge = add_rounded_rect(slide, badge_left, badge_top, badge_size, badge_size, step["accent"])
    badge.text_frame.clear()
    p = badge.text_frame.paragraphs[0]
    p.text = step["num"]
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = "Tahoma"

    title_color = WHITE if is_final else TEXT
    desc_color = RGBColor(0xBB, 0xDE, 0xFB) if is_final else MUTED

    add_textbox(
        slide, left + Inches(0.15), top + Inches(0.12), width - Inches(0.65), Inches(0.35),
        step["title"], font_size=13, bold=True, color=title_color, align=PP_ALIGN.RIGHT,
    )
    add_textbox(
        slide, left + Inches(0.15), top + Inches(0.48), width - Inches(0.2), Inches(0.42),
        step["desc"], font_size=10, color=desc_color, align=PP_ALIGN.RIGHT,
    )

    if is_filter:
        add_textbox(
            slide, left + Inches(0.15), top + height - Inches(0.32), width - Inches(0.2), Inches(0.25),
            "✓  CR < 0.1", font_size=11, bold=True, color=GREEN_DARK, align=PP_ALIGN.CENTER, rtl=False,
        )

    return box


def add_reject_branch(slide, filter_left, filter_top, filter_width):
    """Small side branch showing rejected matrices."""
    rx = filter_left - Inches(1.55)
    ry = filter_top + Inches(0.35)
    box = add_rounded_rect(slide, rx, ry, Inches(1.4), Inches(0.55), RED_BG, RED_SOFT)
    add_textbox(
        slide, rx + Inches(0.05), ry + Inches(0.05), Inches(1.3), Inches(0.45),
        "CR ≥ 0.1\nبازنگری / حذف", font_size=9, bold=True, color=RED_SOFT, align=PP_ALIGN.CENTER,
    )
    # dashed connector from filter to reject box
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        filter_left, filter_top + Inches(0.6),
        rx + Inches(1.4), ry + Inches(0.28),
    )
    conn.line.color.rgb = RED_SOFT
    conn.line.width = Pt(1.5)
    conn.line.dash_style = 3  # dash


def build_slide(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    fill_shape(bg, BG)
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

    # Top bar
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.12))
    fill_shape(top_bar, NAVY)

    # Title
    add_textbox(
        slide, Inches(0.5), Inches(0.28), Inches(9), Inches(0.55),
        "فرآیند گردآوری و تجمیع نظرات خبرگان (AHP)",
        font_size=28, bold=True, color=NAVY,
    )
    add_textbox(
        slide, Inches(0.5), Inches(0.82), Inches(9), Inches(0.35),
        "از جمع‌آوری ماتریس‌های زوجی تا محاسبه وزن نهایی معیارها",
        font_size=13, color=MUTED,
    )

    # Layout: 3 columns × 2 rows (RTL flow: right→left, top→bottom)
    box_w = Inches(2.85)
    box_h = Inches(0.95)
    gap_x = Inches(0.35)
    gap_y = Inches(0.55)
    start_x = Inches(6.55)  # rightmost column (RTL)
    start_y = Inches(1.35)

    positions = []
    for row in range(2):
        for col in range(3):
            idx = row * 3 + col
            if idx >= len(STEPS):
                break
            x = start_x - col * (box_w + gap_x)
            y = start_y + row * (box_h + gap_y)
            positions.append((x, y))

    # Draw steps
    for i, (step, (x, y)) in enumerate(zip(STEPS, positions)):
        add_step_box(slide, x, y, box_w, box_h, step)

    # Arrows between steps (RTL snake flow)
    arrow_pairs = [(0, 1), (1, 2), (2, 3), (3, 4), (4, 5)]
    for a, b in arrow_pairs:
        ax, ay = positions[a]
        bx, by = positions[b]
        if ay == by:  # same row: horizontal arrow
            y_mid = ay + box_h / 2
            if ax > bx:  # right to left
                x1, x2 = ax, bx + box_w
            else:
                x1, x2 = ax + box_w, bx
            conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y_mid, x2, y_mid)
            conn.line.color.rgb = ACCENT
            conn.line.width = Pt(2.5)
            ln = conn.line._ln
            ln.append(parse_xml(
                '<a:tailEnd xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" type="triangle"/>'
            ))
        else:  # row change: vertical then horizontal
            x_mid = ax + box_w / 2
            y1 = ay + box_h
            y2 = by
            conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x_mid, y1, x_mid, y2)
            conn.line.color.rgb = ACCENT
            conn.line.width = Pt(2.5)
            ln = conn.line._ln
            ln.append(parse_xml(
                '<a:tailEnd xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" type="triangle"/>'
            ))

    # Reject branch on step 4
    fx, fy = positions[3]
    add_reject_branch(slide, fx, fy, box_w)

    # Legend panel (left side)
    legend = add_rounded_rect(slide, Inches(0.4), Inches(1.35), Inches(2.5), Inches(2.55), WHITE, BORDER)
    add_textbox(
        slide, Inches(0.55), Inches(1.5), Inches(2.2), Inches(0.3),
        "راهنمای نمادها", font_size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER,
    )

    legends = [
        (NAVY_LIGHT, "مرحله پردازش"),
        (ACCENT, "بررسی سازگاری"),
        (GREEN, "فیلتر پذیرش (CR)"),
        (NAVY, "خروجی نهایی"),
        (RED_SOFT, "رد / بازنگری"),
    ]
    for i, (color, label) in enumerate(legends):
        ly = Inches(1.9) + i * Inches(0.38)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.65), ly, Inches(0.18), Inches(0.18))
        fill_shape(dot, color)
        add_textbox(slide, Inches(0.95), ly - Inches(0.02), Inches(1.8), Inches(0.25), label, font_size=10, color=TEXT)

    # Saaty scale mini reference
    scale_box = add_rounded_rect(slide, Inches(0.4), Inches(4.15), Inches(2.5), Inches(0.85), AMBER_BG, AMBER)
    add_textbox(
        slide, Inches(0.55), Inches(4.28), Inches(2.2), Inches(0.25),
        "مقیاس ساعتی AHP", font_size=11, bold=True, color=AMBER, align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide, Inches(0.55), Inches(4.55), Inches(2.2), Inches(0.3),
        "۱ = برابر  ·  ۳ = متوسط  ·  ۵ = قوی  ·  ۹ = مطلق",
        font_size=9, color=TEXT, align=PP_ALIGN.CENTER,
    )

    # Footer
    add_textbox(
        slide, Inches(0.4), Inches(5.1), Inches(9.2), Inches(0.25),
        "AHP — Analytic Hierarchy Process  |  Saaty",
        font_size=9, color=MUTED, align=PP_ALIGN.LEFT, rtl=False,
    )


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    build_slide(prs)
    output = "/workspace/presentations/فرآیند_تجمیع_نظرات_خبرگان_AHP.pptx"
    prs.save(output)
    print(f"Saved: {output}")


if __name__ == "__main__":
    main()
