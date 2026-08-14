#!/usr/bin/env python3
"""Generate AHP consistency slide — chain + decision box layout."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

BLUE = RGBColor(0x25, 0x63, 0xEB)
BLUE_DARK = RGBColor(0x1D, 0x4E, 0xD8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG = RGBColor(0xFA, 0xFB, 0xFC)
TEXT = RGBColor(0x1E, 0x29, 0x3B)
MUTED = RGBColor(0x64, 0x74, 0x8B)
BORDER = RGBColor(0xE2, 0xE8, 0xF0)
ARROW = RGBColor(0xCB, 0xD5, 0xE1)
GREEN = RGBColor(0x16, 0xA3, 0x34)
GREEN_BG = RGBColor(0xF0, 0xFD, 0xF4)
GREEN_DARK = RGBColor(0x15, 0x80, 0x3D)
RED = RGBColor(0xDC, 0x26, 0x26)
RED_BG = RGBColor(0xFE, 0xF2, 0xF2)
LIGHT_BLUE = RGBColor(0xEF, 0xF6, 0xFF)
ACCENT_BG = RGBColor(0xDB, 0xEA, 0xFE)

CHAIN = [
    {
        "symbol": "λmax",
        "label": "مقدار ویژه بیشینه",
        "formula": "از ماتریس تجمیع‌شده و بردار وزن",
    },
    {
        "symbol": "CI",
        "label": "شاخص سازگاری",
        "formula": "CI = (λmax − n) / (n − 1)",
    },
    {
        "symbol": "CR",
        "label": "نرخ سازگاری",
        "formula": "CR = CI / RI",
        "highlight": True,
    },
]

SIDEBAR_ICONS = 5
ACTIVE_ICON = 2


def set_rtl(paragraph):
    p_pr = paragraph._p.get_or_add_pPr()
    p_pr.set(qn("a:rtl"), "1")


def fill_shape(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_text(slide, left, top, w, h, text, *, size=12, bold=False, color=TEXT,
             align=PP_ALIGN.RIGHT, rtl=True):
    box = slide.shapes.add_textbox(left, top, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    if rtl:
        set_rtl(p)
    run = p.runs[0]
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Tahoma"
    return box


def add_rounded(slide, left, top, w, h, fill, border=None, bw=Pt(1.5)):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = bw
    else:
        s.line.fill.background()
    if s.adjustments:
        s.adjustments[0] = 0.1
    return s


def add_arrow_left(slide, x_right, x_left, y, h=Inches(0.16)):
    w = x_right - x_left
    if w < Inches(0.05):
        return
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x_left, y - h / 2, w, h)
    a.fill.solid()
    a.fill.fore_color.rgb = ARROW
    a.line.fill.background()
    a.rotation = 180.0


def draw_sidebar(slide, slide_w, slide_h):
    sw = Inches(0.55)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, slide_w - sw, 0, sw, slide_h)
    fill_shape(bar, BLUE)

    cx = slide_w - sw / 2
    ys = [Inches(0.9), Inches(1.7), Inches(2.55), Inches(3.4), Inches(4.25)]
    for i, y in enumerate(ys):
        active = i == ACTIVE_ICON
        if active:
            ring = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - Inches(0.22), y - Inches(0.22),
                                          Inches(0.44), Inches(0.44))
            fill_shape(ring, WHITE)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - Inches(0.1), y - Inches(0.1),
                                     Inches(0.2), Inches(0.2))
        fill_shape(dot, BLUE if active else RGBColor(0xDB, 0xEA, 0xFE))


def add_chain_panel(slide, left, top, width, height):
    panel = add_rounded(slide, left, top, width, height, WHITE, BORDER)

    add_text(slide, left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.3),
             "زنجیره محاسبه", size=14, bold=True, color=BLUE_DARK, align=PP_ALIGN.RIGHT)

    # Chain boxes — RTL
    box_w = Inches(1.55)
    box_h = Inches(1.35)
    chain_top = top + Inches(0.55)
    gap = Inches(0.38)
    total_chain = len(CHAIN) * box_w + (len(CHAIN) - 1) * gap
    start_right = left + width - Inches(0.25)

    positions = []
    right = start_right
    for _ in CHAIN:
        l = right - box_w
        positions.append(l)
        right = l - gap

    arrow_y = chain_top + box_h / 2
    for i in range(len(positions) - 1):
        add_arrow_left(slide, positions[i], positions[i + 1] + box_w, arrow_y)

    for (l, step) in zip(positions, CHAIN):
        hl = step.get("highlight", False)
        fill = ACCENT_BG if hl else LIGHT_BLUE
        border = BLUE if hl else BORDER
        add_rounded(slide, l, chain_top, box_w, box_h, fill, border, Pt(2 if hl else 1))

        add_text(slide, l + Inches(0.08), chain_top + Inches(0.12), box_w - Inches(0.16), Inches(0.3),
                 step["symbol"], size=18, bold=True, color=BLUE_DARK if hl else BLUE, align=PP_ALIGN.CENTER, rtl=False)
        add_text(slide, l + Inches(0.08), chain_top + Inches(0.42), box_w - Inches(0.16), Inches(0.28),
                 step["label"], size=10, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
        add_text(slide, l + Inches(0.08), chain_top + Inches(0.72), box_w - Inches(0.16), Inches(0.5),
                 step["formula"], size=8.5, color=MUTED, align=PP_ALIGN.CENTER)

    # n and RI badges
    badge_y = chain_top + box_h + Inches(0.22)
    b1 = add_rounded(slide, left + Inches(0.25), badge_y, Inches(1.3), Inches(0.32), LIGHT_BLUE, BLUE)
    b1.text_frame.clear()
    p = b1.text_frame.paragraphs[0]
    p.text = "n = 8 معیار"
    p.alignment = PP_ALIGN.CENTER
    set_rtl(p)
    run = p.runs[0]
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = BLUE_DARK
    run.font.name = "Tahoma"

    b2 = add_rounded(slide, left + Inches(1.65), badge_y, Inches(1.55), Inches(0.32), LIGHT_BLUE, BLUE)
    b2.text_frame.clear()
    p = b2.text_frame.paragraphs[0]
    p.text = "RI = 1.41 (جدول ساعتی)"
    p.alignment = PP_ALIGN.CENTER
    set_rtl(p)
    run = p.runs[0]
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = BLUE_DARK
    run.font.name = "Tahoma"

    # Example
    add_text(slide, left + Inches(0.25), badge_y + Inches(0.42), width - Inches(0.5), Inches(0.35),
             "مثال: CI = 0.05  →  CR = 0.05 / 1.41 ≈ 0.035  →  قابل قبول ✓",
             size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def add_decision_panel(slide, left, top, width, height):
    panel = add_rounded(slide, left, top, width, height, WHITE, BORDER)

    add_text(slide, left + Inches(0.1), top + Inches(0.15), width - Inches(0.2), Inches(0.3),
             "معیار پذیرش", size=14, bold=True, color=BLUE_DARK, align=PP_ALIGN.CENTER)

    # Large CR threshold
    cr_box = add_rounded(slide, left + Inches(0.15), top + Inches(0.55), width - Inches(0.3), Inches(0.65), BLUE)
    add_text(slide, left + Inches(0.15), top + Inches(0.68), width - Inches(0.3), Inches(0.4),
             "CR < 0.1", size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER, rtl=False)

    # Accept
    ok = add_rounded(slide, left + Inches(0.15), top + Inches(1.35), width - Inches(0.3), Inches(0.55), GREEN_BG, GREEN)
    add_text(slide, left + Inches(0.15), top + Inches(1.48), width - Inches(0.3), Inches(0.3),
             "✓  قابل قبول", size=12, bold=True, color=GREEN_DARK, align=PP_ALIGN.CENTER)

    # Reject
    bad = add_rounded(slide, left + Inches(0.15), top + Inches(2.0), width - Inches(0.3), Inches(0.55), RED_BG, RED)
    add_text(slide, left + Inches(0.15), top + Inches(2.13), width - Inches(0.3), Inches(0.3),
             "✗  نیاز به بازنگری", size=12, bold=True, color=RED, align=PP_ALIGN.CENTER)

    add_text(slide, left + Inches(0.1), top + Inches(2.7), width - Inches(0.2), Inches(0.45),
             "در صورت عدم سازگاری، مقایسه‌های زوجی بازبینی می‌شوند.",
             size=8.5, color=MUTED, align=PP_ALIGN.CENTER)


def add_cr_scale(slide, left, top, width):
    add_text(slide, left, top, width, Inches(0.25),
             "مقیاس نرخ سازگاری", size=10, bold=True, color=TEXT, align=PP_ALIGN.RIGHT)

    bar_y = top + Inches(0.3)
    bar_h = Inches(0.18)
    total_w = width

    # Green zone (0 to 0.1)
    threshold_ratio = 0.35
    green_w = total_w * threshold_ratio
    red_w = total_w - green_w

    green = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, bar_y, green_w, bar_h)
    green.fill.solid()
    green.fill.fore_color.rgb = GREEN_BG
    green.line.color.rgb = GREEN
    green.line.width = Pt(1)

    red_left = left + green_w
    red = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, red_left, bar_y, red_w, bar_h)
    red.fill.solid()
    red.fill.fore_color.rgb = RED_BG
    red.line.color.rgb = RED
    red.line.width = Pt(1)

    # Threshold marker
    marker_x = left + green_w
    marker = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, marker_x - Pt(1.5), bar_y - Inches(0.04), Pt(3), bar_h + Inches(0.08))
    fill_shape(marker, BLUE_DARK)

    add_text(slide, left, bar_y + Inches(0.22), green_w, Inches(0.22),
             "قابل قبول", size=8, color=GREEN_DARK, align=PP_ALIGN.CENTER)
    add_text(slide, red_left, bar_y + Inches(0.22), red_w, Inches(0.22),
             "بازنگری", size=8, color=RED, align=PP_ALIGN.CENTER)
    add_text(slide, marker_x - Inches(0.25), bar_y + Inches(0.22), Inches(0.5), Inches(0.22),
             "0.1", size=8, bold=True, color=BLUE_DARK, align=PP_ALIGN.CENTER, rtl=False)

    add_text(slide, left, bar_y - Inches(0.18), Inches(0.3), Inches(0.18),
             "0", size=8, color=MUTED, align=PP_ALIGN.CENTER, rtl=False)


def build_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    fill_shape(bg, BG)
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

    draw_sidebar(slide, prs.slide_width, prs.slide_height)

    margin = Inches(0.5)
    content_right = prs.slide_width - Inches(0.65)

    add_text(slide, margin, Inches(0.38), content_right - margin, Inches(0.45),
             "بررسی سازگاری ماتریس نهایی", size=24, bold=True, color=TEXT, align=PP_ALIGN.RIGHT)
    add_text(slide, margin, Inches(0.88), content_right - margin, Inches(0.3),
             "پس از محاسبه بردار وزن، سازگاری ماتریس تجمیع‌شده بررسی می‌شود",
             size=11, color=MUTED, align=PP_ALIGN.RIGHT)

    panel_top = Inches(1.35)
    panel_h = Inches(2.85)
    decision_w = Inches(2.2)
    chain_w = content_right - margin - decision_w - Inches(0.2)
    decision_left = content_right - decision_w

    add_chain_panel(slide, margin, panel_top, chain_w, panel_h)
    add_decision_panel(slide, decision_left, panel_top, decision_w, panel_h)

    add_cr_scale(slide, margin, Inches(4.35), content_right - margin)

    # Footer
    add_text(slide, margin, Inches(4.88), Inches(2.0), Inches(0.3),
             "روش پیشنهادی", size=11, bold=True, color=BLUE, align=PP_ALIGN.RIGHT)

    nav_y = Inches(4.9)
    for x in (Inches(3.0), Inches(3.75)):
        c = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, nav_y, Inches(0.28), Inches(0.28))
        fill_shape(c, BLUE)
    add_text(slide, Inches(3.35), nav_y - Inches(0.02), Inches(0.35), Inches(0.3),
             "25", size=11, bold=True, color=TEXT, align=PP_ALIGN.CENTER, rtl=False)


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    build_slide(prs)
    out = "/workspace/presentations/بررسی_سازگاری_ماتریس_نهایی.pptx"
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == "__main__":
    main()
