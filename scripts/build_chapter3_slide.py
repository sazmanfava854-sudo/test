from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

BG = RGBColor(0xFA, 0xFB, 0xFD)
BLUE = RGBColor(0x25, 0x63, 0xEB)
BLUE_DARK = RGBColor(0x1E, 0x40, 0xAF)
BLUE_LIGHT = RGBColor(0xDB, 0xEA, 0xFE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x0F, 0x17, 0x2A)
MUTED = RGBColor(0x64, 0x74, 0x8B)
ARROW = RGBColor(0x94, 0xA3, 0xB8)
LINE = RGBColor(0xE2, 0xE8, 0xF0)
SHADOW = RGBColor(0xD1, 0xD9, 0xE6)

slide.background.fill.solid()
slide.background.fill.fore_color.rgb = BG


def set_text(tf, text, size=16, bold=False, color=DARK, font="B Nazanin", align=PP_ALIGN.CENTER):
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_rect(slide, left, top, width, height, fill, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_box(slide, left, top, width, height, title, subtitle, latin=False, is_output=False):
    add_rect(slide, left + Emu(36576), top + Emu(36576), width, height, SHADOW)

    fill = BLUE_DARK if is_output else BLUE
    box = add_rect(slide, left, top, width, height, fill, BLUE_DARK, Pt(0.75))

    strip_h = Inches(0.05)
    add_rect(slide, left, top, width, strip_h, RGBColor(0x60, 0xA5, 0xFA))

    title_font = "Times New Roman" if latin else "B Nazanin"
    title_size = 22 if latin else 20
    sub_size = 13 if is_output else 14

    title_tb = slide.shapes.add_textbox(left, top + strip_h + Inches(0.28), width, Inches(0.42))
    set_text(title_tb.text_frame, title, title_size, True, WHITE, title_font)

    sub_tb = slide.shapes.add_textbox(left, top + strip_h + Inches(0.72), width, Inches(0.55))
    set_text(sub_tb.text_frame, subtitle, sub_size, False, BLUE_LIGHT, "B Nazanin")


def add_arrow(slide, left, top, width, height):
    chevron = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
    chevron.fill.solid()
    chevron.fill.fore_color.rgb = ARROW
    chevron.line.fill.background()
    chevron.rotation = 180.0


# Header
add_rect(slide, Inches(11.9), Inches(0.52), Inches(0.07), Inches(0.58), BLUE)
title_tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.45), Inches(11.0), Inches(0.62))
set_text(title_tb.text_frame, "فرایند انتخاب و رتبه‌بندی فناوری‌ها", 30, True, DARK, "B Nazanin", PP_ALIGN.RIGHT)

sub_tb = slide.shapes.add_textbox(Inches(0.8), Inches(1.02), Inches(11.0), Inches(0.35))
set_text(sub_tb.text_frame, "چارچوب چهارمرحله‌ای مدل پیشنهادی", 14, False, MUTED, "B Nazanin", PP_ALIGN.RIGHT)

add_rect(slide, Inches(0.8), Inches(1.48), Inches(11.73), Pt(1), LINE)

# Layout (RTL: right -> left)
BW = Inches(2.25)
BH = Inches(1.75)
OW = Inches(1.45)
BY = Inches(3.0)
AW = Inches(0.38)
AH = Inches(0.36)
AY = BY + (BH - AH) / 2

x_cluster = Inches(10.58)
x_ahp = x_cluster - AW - BW
x_dynamic = x_ahp - AW - BW
x_topsis = x_dynamic - AW - BW
x_output = x_topsis - AW - OW

blocks = [
    (x_cluster, "خوشه‌بندی", "گروه‌بندی فناوری‌های مشابه", False, False),
    (x_ahp, "AHP", "تعیین اهمیت نسبی معیارها", True, False),
    (x_dynamic, "تعدیل پویا", "تطبیق وزن‌ها با سناریوی پروژه", False, False),
    (x_topsis, "TOPSIS", "رتبه‌بندی نهایی در خوشه منتخب", True, False),
]

for left, title, sub, latin, out in blocks:
    add_box(slide, left, BY, BW, BH, title, sub, latin, out)

add_box(
    slide,
    x_output,
    BY + (BH - Inches(1.45)) / 2,
    OW,
    Inches(1.45),
    "فناوری منتخب",
    "خروجی نهایی",
    False,
    True,
)

for ax in [
    x_cluster - AW,
    x_ahp - AW,
    x_dynamic - AW,
    x_topsis - AW,
]:
    add_arrow(slide, ax, AY, AW, AH)

# Step numbers (small circles, top-right of each main block)
steps = [(x_cluster, "۱"), (x_ahp, "۲"), (x_dynamic, "۳"), (x_topsis, "۴")]
for left, num in steps:
    cx = left + BW - Inches(0.28)
    cy = BY + Inches(0.12)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, Inches(0.28), Inches(0.28))
    circle.fill.solid()
    circle.fill.fore_color.rgb = WHITE
    circle.line.color.rgb = BLUE_DARK
    circle.line.width = Pt(0.75)
    num_tb = slide.shapes.add_textbox(cx, cy, Inches(0.28), Inches(0.28))
    set_text(num_tb.text_frame, num, 11, True, BLUE_DARK, "B Nazanin")

# Caption
cap_tb = slide.shapes.add_textbox(Inches(0.8), Inches(5.35), Inches(11.73), Inches(0.4))
set_text(
    cap_tb.text_frame,
    "شکل X. چارچوب فرایندی انتخاب فناوری‌های شبکه سلولی",
    15,
    False,
    MUTED,
    "B Nazanin",
)

out_path = "/workspace/artifacts/chapter3-process-diagram.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
